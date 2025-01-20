import streamlit as st
import xgboost

import shap
import matplotlib.pyplot as plt
import matplotlib
import seaborn as sns
from pptx import Presentation
from pptx.util import Inches
from io import BytesIO

st.title(":blue[Streamlit-shap] :sunglasses:")

default_pos_color = "#ff0051"
default_neg_color = "#008bfb"
# Custom colors
positive_color = "#ca0020"
negative_color = "#92c5de"
def add_figure_to_ppt(ppt_template, path_image, save_path):
    ppt = Presentation(ppt_template)
    # Get slide where image shoud be added
    slide = ppt.slides[4]
    # Get slide dimensions (in points)
    slide_width = ppt.slide_width
    slide_height = ppt.slide_height 
    # Image dimensions (in inches) 
    image_width = Inches(12) 
    image_height = Inches(6)
    # Calculate positions to center the image 
    left = (slide_width - image_width) / 2 
    top = (slide_height - image_height) / 2
    # Addinf image
    pic = slide.shapes.add_picture(path_image, left, top, width=image_width, height=image_height)
    # Save ppt
    # ppt.save(save_path)
    return ppt

def explainability_observation(explainer, X, index, style):
    # plt.style.use(style)
    shap_values_idx = explainer(X.loc[[index]])
    fig = plt.figure()
    shap.waterfall_plot(shap_values_idx[0], show = False)
    # Change the colormap of the artists
    for fc in plt.gcf().get_children():
        # Ignore last Rectangle
        for fcc in fc.get_children()[:-1]:
            if (isinstance(fcc, matplotlib.patches.Rectangle)):
                if (matplotlib.colors.to_hex(fcc.get_facecolor()) == default_pos_color):
                    fcc.set_facecolor(positive_color)
                elif (matplotlib.colors.to_hex(fcc.get_facecolor()) == default_neg_color):
                    fcc.set_color(negative_color)
            elif (isinstance(fcc, plt.Text)):
                if (matplotlib.colors.to_hex(fcc.get_color()) == default_pos_color):
                    fcc.set_color(positive_color)
                elif (matplotlib.colors.to_hex(fcc.get_color()) == default_neg_color):
                    fcc.set_color(negative_color)
    return fig


def overall_explainability(explainer, X, style):

    # Set the style
    sns.set(style="whitegrid")

    # Create the figure
    fig, ax = plt.subplots(figsize=(10, 6))  # Adjust the figure size (width, height in inches)

    # Generate SHAP values
    shap_values = explainer(X)
    shap.plots.bar(shap_values, show=False, ax=ax,)
    # Customize the plot appearance
    for fc in plt.gcf().get_children():
        # Ignore last Rectangle
        for fcc in fc.get_children()[:-1]:
            # for the rectangle
            if (isinstance(fcc, matplotlib.patches.Rectangle)):
                if (matplotlib.colors.to_hex(fcc.get_facecolor()) == default_pos_color):
                    fcc.set_facecolor("#FF5733")
            # for the text
            elif (isinstance(fcc, plt.Text)):
                if (matplotlib.colors.to_hex(fcc.get_color()) == default_pos_color):
                    fcc.set_color("#000000")
               
    
    ax.set_title('SHAP Values', fontsize=16)
    ax.set_xlabel('Feature Importance', fontsize=14)
    ax.set_ylabel('Features', fontsize=14)
    ax.grid(True, which='both', linestyle='--', linewidth=0.5)
    ax.tick_params(axis='x', rotation=45, labelsize=12)
    ax.tick_params(axis='y', labelsize=12)
    return fig


st.sidebar.header('Parameters')
# FIXME: Move model into another script
# train XGBoost model
X, y = shap.datasets.adult(n_points=2000)
model = xgboost.XGBClassifier().fit(X, y)

# compute SHAP values
explainer = shap.Explainer(model, X)
idx = st.selectbox(
    'Select observation',
    X.index)

st.write('You selected:', idx)

st.write("Gives the contribution of each feature to the model's output for a specific prediction. The output is sequential view of how features influence the model's prediction.")
st.pyplot(explainability_observation(explainer, X, idx, "light_background"))

st.write("Average impact of each feature in the whole dataset")
fig = overall_explainability(explainer, X, style="light_background")
st.pyplot(fig)


submit = st.button(label='Generate PowerPoint slides')

if submit != '':
    plt.savefig('temp.jpg', format='jpeg', bbox_inches='tight')
    ppt = add_figure_to_ppt('template.pptx', 'temp.jpg', 'updated_ppt.pptx')
    binary_output = BytesIO()
    ppt.save(binary_output)

    st.download_button(label='Click to download PowerPoint',
                                data=binary_output.getvalue(),
                                file_name='test.pptx')